#!/usr/bin/env python3
"""
Generate PowerPoint presentation from markdown slide script.
"""

from pathlib import Path
import re
from pptx import Presentation

ROOT = Path('/Users/junghualiu/case/a2a/qnlp/discocat_qnlp_analysis')
SCRIPT_PATH = ROOT / 'presentations' / 'quantum_nlp_conference_script.md'
OUTPUT_PATH = ROOT / 'information_society' / 'presentations' / 'quantum_nlp_conference_slides.pptx'

SLIDE_HEADING_RE = re.compile(r'^##\s*Slide\s*(\d+)\s*[–-]\s*(?:Title:\s*)?(.+)$')


def parse_script(text: str):
    slides = []
    current = None
    in_notes = False

    for raw_line in text.splitlines():
        line = raw_line.strip()

        match = SLIDE_HEADING_RE.match(line)
        if match:
            if current:
                slides.append(current)
            slide_num = int(match.group(1))
            title = match.group(2).strip()
            current = {'number': slide_num, 'title': title, 'bullets': [], 'notes': []}
            in_notes = False
            continue

        if current is None:
            continue

        if line.startswith('**Speaker Notes**'):
            in_notes = True
            continue

        if not line:
            in_notes = in_notes
            continue

        if in_notes:
            current['notes'].append(raw_line.strip())
        elif line.startswith('- '):
            current['bullets'].append(line[2:].strip())

    if current:
        slides.append(current)

    return slides


def create_presentation(slides):
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    if slides:
        first = slides[0]
        title_slide = prs.slides.add_slide(title_layout)
        title_slide.shapes.title.text = first['title']
        subtitle = title_slide.placeholders[1]
        subtitle.text = '\n'.join(first['bullets']) if first['bullets'] else ''
        if first['notes']:
            title_slide.notes_slide.notes_text_frame.text = '\n'.join(first['notes'])
        slides = slides[1:]

    for slide in slides:
        ppt_slide = prs.slides.add_slide(bullet_layout)
        ppt_slide.shapes.title.text = slide['title']
        text_frame = ppt_slide.shapes.placeholders[1].text_frame
        text_frame.clear()
        for idx, bullet in enumerate(slide['bullets']):
            if idx == 0:
                text_frame.text = bullet
            else:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
        if not slide['bullets']:
            text_frame.text = ''
        if slide['notes']:
            ppt_slide.notes_slide.notes_text_frame.text = '\n'.join(slide['notes'])

    prs.save(OUTPUT_PATH)


def main():
    script_text = SCRIPT_PATH.read_text(encoding='utf-8')
    slides = parse_script(script_text)
    if not slides:
        raise RuntimeError('No slides parsed from script')
    create_presentation(slides)
    print(f'Generated presentation with {len(slides)} slides at {OUTPUT_PATH}')


if __name__ == '__main__':
    main()
